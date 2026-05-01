# HOW TO RUN:
# 1. Open terminal in this folder
# 2. Run: pip install -r requirements.txt
# 3. Run: python app.py
# 4. Open browser: http://localhost:5000
# NOTE: Python 3.9+ required.

from flask import Flask, request, jsonify, render_template, send_file
from flask_cors import CORS
import fitz                          # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
from google import genai             # NEW 2026 SDK — NOT google.generativeai
import io, os, base64, re, json
import struct, zlib
from urllib import request as urlrequest
from urllib import error as urlerror
from werkzeug.utils import secure_filename
import traceback

app = Flask(__name__)
CORS(app) # For production, this should be restricted
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024   # 50 MB hard limit

ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.jpg', '.jpeg', '.png', '.txt', '.md'}

DEFAULT_LOCAL_LLM_URL = 'http://localhost:11434'
DEFAULT_LOCAL_LLM_MODEL = 'llama3.2-vision'
DEFAULT_GEMINI_MODEL = 'gemini-2.5-flash'
DEFAULT_OPENAI_MODEL = 'gpt-4o-mini'
DEFAULT_ANTHROPIC_MODEL = 'claude-3-5-sonnet-latest'

OCR_WARNING_MARKERS = (
    '[No API key',
    '[Scanned page',
    '[Image found',
    '[LOCAL LLM ERROR',
    '[QUOTA EXCEEDED',
    '[INVALID API KEY',
    '[MODEL NOT FOUND',
    '[OCR ERROR'
)


def get_extraction_warning(text: str) -> str:
    if not text:
        return ''
    if any(marker in text for marker in OCR_WARNING_MARKERS):
        return 'Text extraction completed, but one or more images/scanned pages were not OCR processed. Add an API key or enable a working local vision model for image text.'
    return ''


def resolve_cloud_model(api_key: str, requested_model: str = '') -> str:
    requested_model = (requested_model or '').strip()
    if requested_model:
        return requested_model
    if api_key.startswith('sk-ant'):
        return DEFAULT_ANTHROPIC_MODEL
    if api_key.startswith('sk-'):
        return DEFAULT_OPENAI_MODEL
    return DEFAULT_GEMINI_MODEL

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPER FUNCTION: call_ai_ocr(api_key, image, prompt)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def normalize_local_llm_url(base_url: str) -> str:
    base_url = (base_url or DEFAULT_LOCAL_LLM_URL).strip().rstrip('/')
    if base_url.endswith('/api/generate'):
        return base_url
    return f'{base_url}/api/generate'


def normalize_local_llm_base_url(base_url: str) -> str:
    base_url = (base_url or DEFAULT_LOCAL_LLM_URL).strip().rstrip('/')
    for suffix in ('/api/generate', '/api/tags'):
        if base_url.endswith(suffix):
            return base_url[:-len(suffix)].rstrip('/')
    return base_url


def list_local_llm_models(base_url: str) -> list[str]:
    tags_url = f'{normalize_local_llm_base_url(base_url)}/api/tags'
    req = urlrequest.Request(tags_url, method='GET')
    with urlrequest.urlopen(req, timeout=8) as response:
        data = json.loads(response.read().decode('utf-8'))
    models = data.get('models', [])
    names = []
    for model in models:
        name = model.get('name') or model.get('model')
        if name:
            names.append(name)
    return sorted(set(names), key=str.lower)


def call_local_llm_ocr(base_url: str, model: str, image: Image.Image, prompt: str) -> str:
    """Call a local Ollama-compatible vision model for OCR."""
    buffered = io.BytesIO()
    if image.mode in ('RGBA', 'P'):
        image = image.convert('RGB')
    image.save(buffered, format="PNG")
    img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")

    payload = {
        'model': (model or DEFAULT_LOCAL_LLM_MODEL).strip(),
        'prompt': prompt,
        'images': [img_b64],
        'stream': False
    }
    req = urlrequest.Request(
        normalize_local_llm_url(base_url),
        data=json.dumps(payload).encode('utf-8'),
        headers={'Content-Type': 'application/json'},
        method='POST'
    )

    try:
        with urlrequest.urlopen(req, timeout=180) as response:
            result = json.loads(response.read().decode('utf-8'))
    except urlerror.URLError as e:
        return f'[LOCAL LLM ERROR: Could not connect to local LLM at {base_url or DEFAULT_LOCAL_LLM_URL}. {str(e)[:160]}]'

    text = (result.get('response') or '').strip()
    if not text:
        return '[LOCAL LLM ERROR: Empty response. Check that the selected local model supports vision/images.]'
    return text


def call_ai_ocr(
    api_key: str,
    image: Image.Image,
    prompt: str,
    ai_provider: str = 'cloud',
    local_llm_url: str = DEFAULT_LOCAL_LLM_URL,
    local_llm_model: str = DEFAULT_LOCAL_LLM_MODEL,
    cloud_model: str = ''
) -> str:
    """Call an AI Vision model to OCR an image based on the provided API key."""
    if ai_provider == 'local':
        return call_local_llm_ocr(local_llm_url, local_llm_model, image, prompt)

    if not api_key:
        return "[No API key found]"
        
    try:
        if api_key.startswith('sk-ant'):
            import anthropic
            client = anthropic.Anthropic(api_key=api_key)
            # convert image
            buffered = io.BytesIO()
            if image.mode in ('RGBA', 'P'):
                image = image.convert('RGB')
            image.save(buffered, format="PNG")
            img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
            
            response = client.messages.create(
                model=resolve_cloud_model(api_key, cloud_model),
                max_tokens=4000,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}}
                    ]
                }]
            )
            return response.content[0].text.strip()
            
        elif api_key.startswith('sk-'):
            import openai
            client = openai.OpenAI(api_key=api_key)
            buffered = io.BytesIO()
            if image.mode in ('RGBA', 'P'):
                image = image.convert('RGB')
            image.save(buffered, format="PNG")
            img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
            
            response = client.chat.completions.create(
                model=resolve_cloud_model(api_key, cloud_model),
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}}
                    ]
                }]
            )
            return response.choices[0].message.content.strip()
            
        else: # Gemini default
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(
                model=resolve_cloud_model(api_key, cloud_model),
                contents=[prompt, image]
            )
            return response.text.strip()
            
    except Exception as e:
        error_msg = str(e)
        if '429' in error_msg or 'quota' in error_msg.lower():
            return '[QUOTA EXCEEDED — You have used all free requests for this API key. Wait or use a new key.]'
        elif '403' in error_msg or 'API_KEY_INVALID' in error_msg or 'invalid api key' in error_msg.lower():
            return '[INVALID API KEY — Open Settings and check your API key]'
        elif '404' in error_msg or 'not found' in error_msg.lower():
            return '[MODEL NOT FOUND — The requested AI model is unavailable]'
        else:
            return f'[OCR ERROR: {error_msg[:200]}]'

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPER FUNCTION: extract_docx(file_stream) -> str
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def extract_docx(file_stream) -> str:
    doc = Document(file_stream)
    text = []

    # Paragraphs — check for embedded images using XML namespace
    for i, para in enumerate(doc.paragraphs):
        stripped = para.text.strip()
        if stripped:
            text.append(stripped)
        # Detect inline images in paragraph XML
        # Inline images live inside <a:blip> or <pic:pic> tags in the XML
        para_xml = para._element.xml if hasattr(para._element, 'xml') else ''
        if ('graphicData' in para_xml or 'pic:pic' in para_xml or
                'a:blip' in para_xml or 'drawing' in para_xml.lower()):
            text.append('[Image found — content not extracted. Add a Gemini API key to enable image OCR]')

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_cells:
                text.append(' | '.join(row_cells))

    # Shapes / text boxes (inline shapes in body XML)
    try:
        body_xml = doc.element.body.xml if hasattr(doc.element.body, 'xml') else ''
        # Count image references that weren't caught in paragraphs
        blip_count = len(re.findall(r'<a:blip\s', body_xml))
        already_noted = sum(1 for t in text if 'Image found' in t)
        extra = blip_count - already_noted
        for _ in range(max(0, extra)):
            text.append('[Image found — content not extracted. Add a Gemini API key to enable image OCR]')
    except Exception:
        pass  # never crash extraction over image detection

    return '\n'.join(text)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPER FUNCTION: extract_pptx(file_stream) -> str
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def extract_pptx(file_stream) -> str:
    prs = Presentation(file_stream)
    text = []

    for i, slide in enumerate(prs.slides):
        text.append(f'--- Slide {i + 1} ---')
        slide_has_image = False

        for shape in slide.shapes:
            # Text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        text.append(line)

            # Tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        text.append(' | '.join(cells))

            # Detect picture shapes — shape_type 13 is MSO_SHAPE_TYPE.PICTURE
            try:
                from pptx.enum.shapes import PP_PLACEHOLDER
                # shape.shape_type == 13 means it is a picture
                if shape.shape_type == 13:
                    slide_has_image = True
                # Also catch placeholder images (shape_type 14 = PLACEHOLDER with image)
                elif hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    ph = shape.placeholder_format
                    # placeholder type 18 = PICTURE placeholder
                    if str(ph.type) in ('PP_PLACEHOLDER.PICTURE', '18'):
                        slide_has_image = True
            except Exception:
                pass

            # Check shape XML for embedded images as a fallback
            try:
                shape_xml = shape._element.xml if hasattr(shape._element, 'xml') else ''
                if 'p:pic' in shape_xml or 'a:blip' in shape_xml:
                    slide_has_image = True
            except Exception:
                pass

        if slide_has_image:
            text.append(f'[Image on slide {i + 1} — content not extracted. Add a Gemini API key to enable image OCR]')

        # Speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text.append(f'[Notes: {notes}]')

    return '\n'.join(text)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPER FUNCTION: extract_pdf(file_bytes, api_key) -> dict
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def extract_pdf(
    file_bytes: bytes,
    api_key: str,
    ai_provider: str = 'cloud',
    local_llm_url: str = DEFAULT_LOCAL_LLM_URL,
    local_llm_model: str = DEFAULT_LOCAL_LLM_MODEL,
    cloud_model: str = ''
) -> dict:
    has_ai_ocr = bool(api_key) or ai_provider == 'local'
    pdf = fitz.open(stream=file_bytes, filetype='pdf')
    page_count = len(pdf)
    text_parts = [None] * page_count
    methods_used = set()
    
    ocr_tasks = []

    # Fast sequential pass for direct text and pixmap grabbing (PyMuPDF objects aren't thread-safe)
    for page_num in range(page_count):
        page = pdf[page_num]
        direct_text = page.get_text().strip()

        if len(direct_text) > 50:
            text_parts[page_num] = f'Page {page_num + 1}:\n{direct_text}'
            methods_used.add('direct')
        else:
            if has_ai_ocr:
                pix = page.get_pixmap(dpi=200)
                img_bytes = pix.tobytes('png')
                ocr_tasks.append((page_num, img_bytes))
            else:
                text_parts[page_num] = (
                    f'Page {page_num + 1}: [Scanned page — add an AI API key or enable Local LLM '
                    f'in Settings to extract this page via OCR]'
                )
                methods_used.add('skipped')
                
    # Threaded pass for network I/O bound OCR calls to speed up queue
    if ocr_tasks:
        from concurrent.futures import ThreadPoolExecutor
        methods_used.add('ocr')
        
        def process_ocr(task):
            p_num, i_bytes = task
            img = Image.open(io.BytesIO(i_bytes))
            text = call_ai_ocr(
                api_key, img,
                'Extract ALL text from this scanned document page exactly as it appears. '
                'Preserve all tables, lists, headings, and structure. '
                'Return ONLY the raw extracted text — no commentary, no markdown.',
                ai_provider,
                local_llm_url,
                local_llm_model,
                cloud_model
            )
            return p_num, f'Page {p_num + 1} (OCR):\n{text}'
            
        # Using 5 concurrent threads avoids rate-limit spikes but vastly speeds up multipage PDFs
        with ThreadPoolExecutor(max_workers=5) as executor:
            for p_num, result_text in executor.map(process_ocr, ocr_tasks):
                text_parts[p_num] = result_text

    method = 'mixed' if len(methods_used) > 1 else (methods_used.pop() if methods_used else 'direct')
    return {
        'text': '\n\n'.join(filter(None, text_parts)),
        'pages': page_count,
        'method': method
    }

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPER FUNCTION: extract_image(file_bytes, api_key) -> str
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def extract_image(
    file_bytes: bytes,
    api_key: str,
    ai_provider: str = 'cloud',
    local_llm_url: str = DEFAULT_LOCAL_LLM_URL,
    local_llm_model: str = DEFAULT_LOCAL_LLM_MODEL,
    cloud_model: str = ''
) -> str:
    if not api_key and ai_provider != 'local':
        return '[No API key — open Settings to add your key or enable Local LLM to process images]'
    img = Image.open(io.BytesIO(file_bytes))
    # Convert RGBA to RGB to avoid PNG transparency issues
    if img.mode in ('RGBA', 'P'):
        img = img.convert('RGB')
    return call_ai_ocr(
        api_key, img,
        'Extract ALL text from this image exactly as it appears. '
        'Preserve tables, lists, headings and structure. '
        'Return ONLY the raw extracted text — no commentary, no markdown.',
        ai_provider,
        local_llm_url,
        local_llm_model,
        cloud_model
    )

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ROUTE 1: GET /
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ROUTE 2: POST /extract
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
@app.route('/extract', methods=['POST'])
def extract():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({"success": False, "error": "No files or pasted text were uploaded."}), 400

        api_key = (
            request.headers.get('X-AI-API-Key', '') or
            request.headers.get('X-Gemini-API-Key', '')
        ).strip()
        ai_provider = request.headers.get('X-AI-Provider', 'cloud').strip().lower()
        if ai_provider not in ('cloud', 'local'):
            ai_provider = 'cloud'
        local_llm_url = request.headers.get('X-Local-LLM-URL', DEFAULT_LOCAL_LLM_URL).strip()
        local_llm_model = request.headers.get('X-Local-LLM-Model', DEFAULT_LOCAL_LLM_MODEL).strip()
        cloud_model = request.headers.get('X-Cloud-AI-Model', '').strip()
        
        results = []
        combined_text_parts = []
        total_words = 0
        files_processed = 0
        files_failed = 0
        files_partial = 0

        for file in files:
            filename = secure_filename(file.filename)
            _, ext = os.path.splitext(filename)

            if ext.lower() not in ALLOWED_EXTENSIONS:
                results.append({"filename": filename, "error": "File type not allowed."})
                files_failed += 1
                continue

            file_bytes = file.read()
            
            if len(file_bytes) > app.config['MAX_CONTENT_LENGTH']:
                results.append({"filename": filename, "error": "File exceeds 50MB limit."})
                files_failed += 1
                continue

            result = {
                "filename": filename,
                "text": "",
                "pages": 1,
                "method": "direct",
                "word_count": 0,
                "char_count": 0,
                "warning": None,
                "partial": False,
                "error": None
            }

            try:
                if ext.lower() == '.pdf':
                    pdf_data = extract_pdf(file_bytes, api_key, ai_provider, local_llm_url, local_llm_model, cloud_model)
                    result['text'] = pdf_data['text']
                    result['pages'] = pdf_data['pages']
                    result['method'] = pdf_data['method']
                elif ext.lower() == '.docx':
                    result['text'] = extract_docx(io.BytesIO(file_bytes))
                elif ext.lower() == '.pptx':
                    result['text'] = extract_pptx(io.BytesIO(file_bytes))
                elif ext.lower() in ['.txt', '.md']:
                    result['text'] = file_bytes.decode('utf-8', errors='replace')
                    result['method'] = 'paste'
                elif ext.lower() in ['.jpg', '.jpeg', '.png']:
                    result['text'] = extract_image(file_bytes, api_key, ai_provider, local_llm_url, local_llm_model, cloud_model)
                    result['method'] = 'local-llm' if ai_provider == 'local' else 'cloud-ocr'

                # Gracefully process text even if it contains an OCR error or missing key warning
                warning = get_extraction_warning(result['text'])
                if warning:
                    result['warning'] = warning
                    result['partial'] = True
                    files_partial += 1

                result['word_count'] = len(result['text'].split())
                result['char_count'] = len(result['text'])
                combined_text_parts.append(f'═'*60 + f'\n📄 {filename}' + '\n' + '═'*60 + '\n\n' + result['text'])
                total_words += result['word_count']
                files_processed += 1

            except Exception as e:
                print(f"Error processing {filename}:")
                traceback.print_exc()
                result['error'] = f"Server error processing file: {str(e)[:100]}"
                files_failed += 1
            
            results.append(result)

        combined_text = '\n\n'.join(combined_text_parts)
        
        return jsonify({
            "success": True,
            "results": results,
            "combined_text": combined_text,
            "total_words": total_words,
            "total_chars": len(combined_text),
            "total_tokens": len(combined_text) // 4,
            "files_processed": files_processed,
            "files_failed": files_failed,
            "files_partial": files_partial
        })

    except Exception as e:
        traceback.print_exc()
        print(f"Error in /extract: {e}")
        from werkzeug.exceptions import RequestEntityTooLarge
        if isinstance(e, RequestEntityTooLarge):
             return jsonify({"success": False, "error": "One or more files exceeded the 50MB total size limit."}), 413
        return jsonify({"success": False, "error": "An unexpected server error occurred."}), 500

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ROUTE 3: POST /test-api-key
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
@app.route('/test-api-key', methods=['POST'])
def test_api_key():
    try:
        data = request.get_json()
        if not data:
            return jsonify({'valid': False, 'error': 'No data received'}), 400

        key = data.get('api_key', '').strip()
        cloud_model = data.get('cloud_model', '').strip()

        if not key:
            return jsonify({'valid': False, 'error': 'API key is empty'}), 400

        # Basic format check first (fast rejection before hitting API)
        if not (key.startswith('AIza') or key.startswith('sk-')):
            return jsonify({'valid': False, 'error': "Key must start with 'AIza' (Gemini), 'sk-proj' (OpenAI), or 'sk-ant' (Anthropic)."}), 400

        # Real live test — generate a tiny 1x1 white PNG in memory
        def make_1x1_png():
            sig = b'\x89PNG\r\n\x1a\n'
            def chunk(name, data):
                c = struct.pack('>I', len(data)) + name + data
                return c + struct.pack('>I', zlib.crc32(name + data) & 0xffffffff)
            ihdr = chunk(b'IHDR', struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0))
            idat = chunk(b'IDAT', zlib.compress(b'\x00\xff\xff\xff'))
            iend = chunk(b'IEND', b'')
            return sig + ihdr + idat + iend

        test_image = Image.open(io.BytesIO(make_1x1_png()))

        # Route test to correct client
        out_msg = ""
        if key.startswith('sk-ant'):
            import anthropic
            client = anthropic.Anthropic(api_key=key)
            buffered = io.BytesIO()
            test_image.save(buffered, format="PNG")
            img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
            client.messages.create(
                model=resolve_cloud_model(key, cloud_model), max_tokens=10,
                messages=[{"role": "user", "content": [{"type": "text", "text": "Reply OK"}]}]
            )
            out_msg = "✅ API key works! Anthropic responded successfully."
        elif key.startswith('sk-'):
            import openai
            client = openai.OpenAI(api_key=key)
            client.chat.completions.create(
                model=resolve_cloud_model(key, cloud_model),
                messages=[{"role": "user", "content": "Reply OK"}]
            )
            out_msg = "✅ API key works! OpenAI responded successfully."
        else:
            client = genai.Client(api_key=key)
            response = client.models.generate_content(
                model=resolve_cloud_model(key, cloud_model),
                contents=['Reply with the single word: OK', test_image]
            )
            _ = response.text
            out_msg = "✅ API key works! Gemini responded successfully."

        return jsonify({
            'valid': True,
            'message': out_msg
        })

    except Exception as e:
        error_msg = str(e)
        if '429' in error_msg or 'quota' in error_msg.lower() or 'insufficient' in error_msg.lower():
            return jsonify({'valid': False, 'error': '⚠️ Quota exceeded / Insufficient balance — key is valid but hit a limit.'}), 200
        elif '403' in error_msg or 'API_KEY_INVALID' in error_msg or 'invalid' in error_msg.lower() or 'authentication' in error_msg.lower():
            return jsonify({'valid': False, 'error': '❌ Invalid API key — please check you copied it correctly.'}), 200
        elif '404' in error_msg or 'not found' in error_msg.lower():
            return jsonify({'valid': False, 'error': '❌ Model not found — try again in a moment'}), 200
        else:
            return jsonify({'valid': False, 'error': f'❌ Test failed: {error_msg[:120]}'}), 200

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ROUTE 4: POST /download-txt
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
@app.route('/test-local-llm', methods=['POST'])
def test_local_llm():
    try:
        data = request.get_json() or {}
        base_url = data.get('base_url', DEFAULT_LOCAL_LLM_URL).strip()
        model = data.get('model', DEFAULT_LOCAL_LLM_MODEL).strip()

        def make_1x1_png():
            sig = b'\x89PNG\r\n\x1a\n'

            def chunk(name, chunk_data):
                c = struct.pack('>I', len(chunk_data)) + name + chunk_data
                return c + struct.pack('>I', zlib.crc32(name + chunk_data) & 0xffffffff)

            ihdr = chunk(b'IHDR', struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0))
            idat = chunk(b'IDAT', zlib.compress(b'\x00\xff\xff\xff'))
            iend = chunk(b'IEND', b'')
            return sig + ihdr + idat + iend

        test_image = Image.open(io.BytesIO(make_1x1_png()))
        response_text = call_local_llm_ocr(
            base_url,
            model,
            test_image,
            'Reply with the single word OK.'
        )
        if response_text.startswith('[LOCAL LLM ERROR:'):
            return jsonify({'valid': False, 'error': response_text}), 200

        return jsonify({
            'valid': True,
            'message': f'Local LLM responded successfully using {model}.'
        })
    except Exception as e:
        return jsonify({'valid': False, 'error': f'Local LLM test failed: {str(e)[:160]}'}), 200


@app.route('/local-llm-models', methods=['POST'])
def local_llm_models():
    try:
        data = request.get_json() or {}
        base_url = data.get('base_url', DEFAULT_LOCAL_LLM_URL).strip()
        models = list_local_llm_models(base_url)
        return jsonify({
            'success': True,
            'models': models,
            'count': len(models)
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'models': [],
            'error': f'Could not load local models: {str(e)[:160]}'
        }), 200


@app.route('/download-txt', methods=['POST'])
def download_txt():
    data = request.get_json()
    text = data.get('text', '')
    filename = data.get('filename', 'doctotext_output.txt')
    filename = secure_filename(filename)
    if not filename.endswith('.txt'):
        filename += '.txt'
    
    buffer = io.BytesIO(text.encode('utf-8'))
    buffer.seek(0)
    return send_file(
        buffer,
        as_attachment=True,
        download_name=filename,
        mimetype='text/plain; charset=utf-8'
    )

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ROUTE 5: POST /download-docx
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
@app.route('/download-docx', methods=['POST'])
def download_docx():
    data = request.get_json()
    text = data.get('text', '')
    filename = data.get('filename', 'doctotext_output.docx')
    filename = secure_filename(filename)
    if not filename.endswith('.docx'):
        filename += '.docx'
    
    # Strip illegal XML characters that crash python-docx
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)

    doc = Document()
    doc.add_heading('DocToText Export', 0)
    for line in text.split('\n'):
        if line.startswith('═') or line.startswith('---'):
            doc.add_paragraph('─' * 40)
        elif line.strip():
            doc.add_paragraph(line)
            
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return send_file(
        buffer,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ROUTE 6: POST /split-chunks
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
@app.route('/split-chunks', methods=['POST'])
def split_chunks():
    data = request.get_json()
    text = data.get('text', '')
    chunk_size = data.get('chunk_size', 4000)
    chars_per_chunk = chunk_size * 4
    
    chunks = []
    current_pos = 0
    while current_pos < len(text):
        end_pos = current_pos + chars_per_chunk
        if end_pos >= len(text):
            chunks.append(text[current_pos:])
            break
        
        # Find the last space to avoid splitting mid-word
        last_space = text.rfind(' ', current_pos, end_pos)
        if last_space == -1 or last_space <= current_pos:
             # No space found, just split at the character limit
            split_at = end_pos
        else:
            split_at = last_space
            
        chunks.append(text[current_pos:split_at].rstrip())
        current_pos = split_at + 1

    return jsonify({
        "chunks": chunks,
        "total_chunks": len(chunks),
        "chunk_size_tokens": chunk_size,
        "chunk_size_chars": chars_per_chunk
    })

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# APP ENTRY POINT:
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
